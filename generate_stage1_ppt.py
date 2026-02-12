
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some brand colors (Modern Tech vibe: Dark Blue, Teal, White)
    PRIMARY_COLOR = RGBColor(0, 51, 102)    # Dark Blue
    ACCENT_COLOR = RGBColor(0, 153, 153)    # Teal
    TEXT_COLOR = RGBColor(50, 50, 50)       # Dark Grey

    def add_slide(layout_index, title_text, content_text=None, subtitle_text=None):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.bold = True

        # Set Content/Subtitle
        if content_text:
            if layout_index == 1: # Title Slide often has subtitle placeholder at index 1
                 subtitle = slide.placeholders[1]
                 subtitle.text = content_text
            else:
                # Content placeholder usually index 1
                content = slide.placeholders[1]
                content.text = content_text


    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Decyphr\nAutomated Deep Data Analysis & Visualization Toolkit"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Capstone Project - Stage 1 Evaluation\n\nPresented By:\n[Group Members]\n\nSupervisor: [Name]\nJanuary 2026"


    # --- Slide 2: Context / The "Hook" ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Context & Motivation"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear() 

    p = text_frame.paragraphs[0]
    p.text = "\"Data is abundant, but insights are expensive.\""
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR

    p = text_frame.add_paragraph()
    p.text = "The Reality: "
    p.font.bold = True
    p.text += "Modern organizations are drowning in data."
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "The Statistic: "
    p.font.bold = True
    p.text += "Data Scientists spend 60%–80% of their time just cleaning and organizing data (CrowdFlower, Forbes)."
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "The Need: "
    p.font.bold = True
    p.text += "A bridge between raw data and decision-ready insights."
    p.level = 0


    # --- Slide 3: Problem Statement (Clarity - 10%) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.text = "The Core Issues with Traditional EDA:"

    bullets = [
        ("Inefficiency & Repetition", "Writing boilerplate code (Pandas, Matplotlib) for every new dataset exploits high-value human capital."),
        ("Inconsistency", "Quality varies widely; critical steps (normality tests, missing value analysis) often skipped due to time pressure."),
        ("Superficiality", "Manual analysis often sticks to basic plots, missing complex, non-linear relationships."),
        ("Integration Gap", "Analysis scattered across notebooks, hard to share with non-technical stakeholders.")
    ]

    for bold_text, normal_text in bullets:
        p = text_frame.add_paragraph()
        p.text = f"{bold_text}: "
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        run = p.add_run()
        run.text = normal_text
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR
        p.level = 0


    # --- Slide 4: Solution Overview ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Solution: Decyphr"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "Concept: "
    p.font.bold = True
    p.text += "Intelligent toolkit automating the end-to-end EDA pipeline."
    
    p = text_frame.add_paragraph()
    p.text = "Mechanism: "
    p.font.bold = True
    p.text += "\"Single Line of Code\" execution."

    p = text_frame.add_paragraph()
    p.text = "Differentiation: "
    p.font.bold = True
    p.text += "Goes beyond basic stats to include ML-driven insights (outlier detection, predictive power scores)."


    # --- Slide 5: Significance (10%) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Significance of the Problem"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.text = "Why Solving This Matters:"

    significance_points = [
        ("Productivity (Time = Money)", "Reduces time-to-insight from days to minutes."),
        ("Reliability & Standardization", "Enforces a rigorous 'Standard of Care'. No check is forgotten."),
        ("Democratization", "Produces presentation-ready reports for non-technical stakeholders."),
        ("Foundational Quality", "Better EDA = Better ML Models.")
    ]

    for bold_text, normal_text in significance_points:
        p = text_frame.add_paragraph()
        p.text = f"{bold_text}: "
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        run = p.add_run()
        run.text = normal_text
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR


    # --- Slide 6: Basis for the Problem ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Basis for the Problem"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "Industry Validation:"
    p.font.bold = True
    
    p = text_frame.add_paragraph()
    p.text = "• The rise of AutoML necessitates AutoEDA."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "• Existing tools (Pandas-Profiling) often crash on big data or lack deep analysis."
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "References:"
    p.font.bold = True
    p.space_before = Pt(20)

    refs = [
        "Anaconda regarding State of Data Science",
        "Forbes on Data Cleaning time costs",
        "Tukey, J. W. (1977). Exploratory Data Analysis"
    ]
    for ref in refs:
         p = text_frame.add_paragraph()
         p.text = f"• {ref}"
         p.level = 1
         p.font.size = Pt(14)


    # --- Slide 7: High Level Objective ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "High-Level Objective"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    
    objectives = [
        "Ingest heterogeneous data sources (CSV, SQL, JSON) automatically.",
        "Detect data types and apply appropriate statistical, ML, and heuristic tests.",
        "Visualize complex relationships interactively without user code.",
        "Export self-contained, shareable reports (HTML/PDF)."
    ]
    
    for obj in objectives:
        p = text_frame.add_paragraph()
        p.text = f"• {obj}"
        p.font.size = Pt(20)


    prs.save('Decyphr_Stage1_Presentation.pptx')
    print("Presentation saved successfully as Decyphr_Stage1_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
